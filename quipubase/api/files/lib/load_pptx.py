import typing as tp
from dataclasses import dataclass

import base64c as base64
import typing_extensions as tpe
from pptx import Presentation

from ._base import Artifact


@dataclass
class PptxLoader(Artifact):
    ref: tpe.Annotated[
        str,
        tpe.Doc(
            """
            This `ref` can represent one out of three things:

            - An HTTP URL.
            - A file path (temporary or not) within the local filesystem.
            - A text file content.
            """
        ),
    ]

    def extract(self) -> tp.Generator[str, None, None]:
        """
        Extract text and images from PowerPoint presentations.

        Yields:
            Text content and image data as HTML
        """
        file_path = self.retrieve()

        try:
            # Load the presentation
            prs = Presentation(file_path.as_posix())

            # Extract content from each slide
            for slide_index, slide in enumerate(prs.slides):
                # Mark the slide number
                yield f"<h2>Slide {slide_index + 1}</h2>"

                # Process each shape in the slide
                for shape in slide.shapes:
                    # Extract text from text frames
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        if text_frame.text:
                            # Build paragraph HTML
                            parts = []
                            for paragraph in text_frame.paragraphs:
                                if paragraph.text:
                                    parts.append(f"<p>{paragraph.text}</p>")
                            if parts:
                                yield "\n".join(parts)

                    # Extract images (shape type 13 is a picture)
                    if hasattr(shape, "shape_type") and shape.shape_type == 13:
                        try:
                            image = shape.image
                            image_data = base64.b64encode(image.blob).decode()
                            yield f'<img style="width: 24em;" src="data:image/png;base64,{image_data}" />'
                        except Exception as img_error:
                            yield f"<p>Error extracting image: {str(img_error)}</p>"

                    # Handle tables
                    if hasattr(shape, "has_table") and shape.has_table:
                        try:
                            table_html = [
                                "<table border='1' style='border-collapse: collapse;'>"
                            ]
                            for row in shape.table.rows:
                                table_html.append("<tr>")
                                for cell in row.cells:
                                    if cell.text:
                                        table_html.append(f"<td>{cell.text}</td>")
                                    else:
                                        table_html.append("<td></td>")
                                table_html.append("</tr>")
                            table_html.append("</table>")
                            yield "\n".join(table_html)
                        except Exception as table_error:
                            yield f"<p>Error extracting table: {str(table_error)}</p>"

        except Exception as e:
            # Handle errors gracefully
            yield f"Error processing PowerPoint file: {str(e)}"
